from __future__ import annotations

import json
import shutil
import uuid
from datetime import datetime
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent
try:
    from dotenv import load_dotenv
    load_dotenv(BASE_DIR / ".env")
except ImportError:
    pass

from typing import Dict, List, Optional

from fastapi import BackgroundTasks, FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse

from speech import audio, clip, score, tts
from speech.align import align_words
from speech.coach import (
    get_answers_for_questions,
    get_answers_for_text,
    get_audience_questions,
    get_gemini_tips,
    get_interview_analysis,
    get_interview_questions,
    get_text_script_feedback_questions,
)
from speech.phoneme_recognizer import recognize_phonemes
from speech.transcribe import transcribe

DATA_DIR = BASE_DIR / "data"
UPLOADS_DIR = DATA_DIR / "uploads"
JOBS_DIR = DATA_DIR / "jobs"
CLIPS_DIR = DATA_DIR / "clips"
TTS_DIR = DATA_DIR / "tts"
SAMPLES_DIR = DATA_DIR / "samples"
PRESENTATION_DIR = DATA_DIR / "presentations"
TEXT_DIR = DATA_DIR / "text"
INTERVIEW_DIR = DATA_DIR / "interview"

for path in [UPLOADS_DIR, JOBS_DIR, CLIPS_DIR, TTS_DIR, SAMPLES_DIR, PRESENTATION_DIR, TEXT_DIR, INTERVIEW_DIR]:
    path.mkdir(parents=True, exist_ok=True)

app = FastAPI(title="AccentCoach API", version="0.1.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://127.0.0.1:3000"],
    allow_origin_regex=r"^http://(localhost|127\.0\.0\.1)(:\d+)?$",
    allow_methods=["*"],
    allow_headers=["*"],
    allow_credentials=True,
)


@app.get("/api/health")
def health() -> Dict[str, bool]:
    return {"ok": True}


def _job_file(job_id: str) -> Path:
    return JOBS_DIR / f"{job_id}.json"


def _write_job(job_id: str, payload: Dict) -> None:
    payload["job_id"] = job_id
    with _job_file(job_id).open("w", encoding="utf-8") as fp:
        json.dump(payload, fp, indent=2)


def _read_job(job_id: str) -> Dict:
    path = _job_file(job_id)
    if not path.exists():
        raise HTTPException(status_code=404, detail="Job not found")
    with path.open("r", encoding="utf-8") as fp:
        return json.load(fp)


@app.post("/api/analyze")
async def analyze(
    background_tasks: BackgroundTasks,
    file: Optional[UploadFile] = File(default=None),
    sample: Optional[str] = None,
):
    if not file and not sample:
        raise HTTPException(status_code=400, detail="Upload an audio file or set sample=true")

    job_id = uuid.uuid4().hex
    upload_name = f"{job_id}_{file.filename}" if file else f"{job_id}_sample.wav"
    upload_path = UPLOADS_DIR / upload_name

    if file:
        content = await file.read()
        with open(upload_path, "wb") as out:
            out.write(content)
    else:
        sample_file = _select_sample(sample)
        shutil.copy(sample_file, upload_path)

    job_record = {
        "status": "processing",
        "created_at": datetime.utcnow().isoformat() + "Z",
    }
    _write_job(job_id, job_record)

    background_tasks.add_task(_process_job, job_id, upload_path)
    return {"job_id": job_id}


@app.get("/api/job/{job_id}")
def get_job(job_id: str):
    return _read_job(job_id)


@app.post("/api/presentation/analyze")
async def analyze_presentation(
    background_tasks: BackgroundTasks,
    slides: UploadFile | None = File(default=None),
    video: UploadFile | None = File(default=None),
    audio: UploadFile | None = File(default=None),
):
    if not slides or (not video and not audio):
        raise HTTPException(status_code=400, detail="Upload slides and a presentation video or audio recording.")

    job_id = f"presentation_{uuid.uuid4().hex}"
    job_record = {
        "status": "processing",
        "created_at": datetime.utcnow().isoformat() + "Z",
        "type": "presentation",
    }
    _write_job(job_id, job_record)

    job_dir = PRESENTATION_DIR / job_id
    job_dir.mkdir(parents=True, exist_ok=True)
    slides_path = job_dir / slides.filename
    media_file = video or audio
    media_path = job_dir / media_file.filename

    slides_content = await slides.read()
    media_content = await media_file.read()
    with slides_path.open("wb") as fp:
        fp.write(slides_content)
    with media_path.open("wb") as fp:
        fp.write(media_content)

    background_tasks.add_task(_process_presentation_job, job_id, slides_path, media_path)
    return {"job_id": job_id}


@app.get("/api/presentation/job/{job_id}")
def get_presentation_job(job_id: str):
    return _read_job(job_id)


@app.post("/api/presentation/job/{job_id}/answers")
def get_presentation_answers(job_id: str) -> Dict:
    """Generate Gemini-powered answers for the audience questions in this job."""
    job = _read_job(job_id)
    if job.get("status") != "done" or not job.get("result"):
        raise HTTPException(status_code=400, detail="Job not ready or no result")
    result = job["result"]
    questions = result.get("audience_questions") or []
    if not questions:
        return {"answers": []}

    transcript = ""
    slide_texts: List[str] = []
    if result.get("analysis"):
        ana = result["analysis"]
        if ana.get("transcript", {}).get("text"):
            transcript = ana["transcript"]["text"]
        if ana.get("slides_text"):
            slide_texts = ana["slides_text"]

    answers = get_answers_for_questions(transcript=transcript, slide_texts=slide_texts, questions=questions)
    return {"answers": answers}


@app.post("/api/text/analyze")
async def analyze_text(
    background_tasks: BackgroundTasks,
    text: Optional[str] = Form(default=None),
    file: Optional[UploadFile] = File(default=None),
):
    """Analyze text or uploaded document. Returns script, feedback, and questions. No video required."""
    content = ""
    if text and text.strip():
        content = text.strip()
    elif file:
        if file.filename:
            suffix = Path(file.filename).suffix.lower()
            if suffix not in {".txt", ".pdf", ".docx"}:
                raise HTTPException(
                    status_code=400,
                    detail="Upload a TXT, PDF, or DOCX file.",
                )
        file_content = await file.read()
        temp_dir = TEXT_DIR / "uploads"
        temp_dir.mkdir(parents=True, exist_ok=True)
        upload_path = temp_dir / f"{uuid.uuid4().hex}_{file.filename or 'upload'}"
        with upload_path.open("wb") as fp:
            fp.write(file_content)
        content = _extract_text_from_file(upload_path)
    else:
        raise HTTPException(status_code=400, detail="Provide text or upload a document (TXT, PDF, DOCX).")

    if not content:
        raise HTTPException(status_code=400, detail="Could not extract text. Try a different file or paste text.")

    job_id = f"text_{uuid.uuid4().hex}"
    job_record = {
        "status": "processing",
        "created_at": datetime.utcnow().isoformat() + "Z",
        "type": "text",
    }
    _write_job(job_id, job_record)
    background_tasks.add_task(_process_text_job, job_id, content)
    return {"job_id": job_id}


@app.get("/api/text/job/{job_id}")
def get_text_job(job_id: str):
    return _read_job(job_id)


@app.post("/api/text/job/{job_id}/answers")
def get_text_answers(job_id: str) -> Dict:
    """Generate answers for the questions in this text job."""
    job = _read_job(job_id)
    if job.get("status") != "done" or not job.get("result"):
        raise HTTPException(status_code=400, detail="Job not ready or no result")
    result = job["result"]
    questions = result.get("questions") or []
    content = result.get("content", "")
    if not questions:
        return {"answers": []}
    answers = get_answers_for_text(content=content, questions=questions)
    return {"answers": answers}


def _process_text_job(job_id: str, content: str) -> None:
    job = _read_job(job_id)
    try:
        data = get_text_script_feedback_questions(content)
        if not data:
            job["status"] = "error"
            job["error"] = "Failed to generate script and feedback. Please try again."
            return
        job.update(
            {
                "status": "done",
                "completed_at": datetime.utcnow().isoformat() + "Z",
                "result": {
                    "content": content[:5000],
                    "script": data.get("script", ""),
                    "feedback": data.get("feedback", []),
                    "questions": data.get("questions", []),
                },
            }
        )
    except Exception as exc:
        job["status"] = "error"
        job["error"] = str(exc)
    finally:
        _write_job(job_id, job)


@app.post("/api/interview/start")
async def start_interview(
    company_name: str = Form(...),
    job_position: str = Form(...),
    company_mission: str = Form(default=""),
    qualifications: str = Form(default=""),
    resume: Optional[UploadFile] = File(default=None),
):
    """Create interview job and generate questions. Qualifications from form or resume."""
    quals = qualifications.strip()
    if resume and resume.filename:
        suffix = Path(resume.filename).suffix.lower()
        if suffix in {".txt", ".pdf", ".docx"}:
            content = await resume.read()
            temp = INTERVIEW_DIR / "uploads"
            temp.mkdir(parents=True, exist_ok=True)
            path = temp / f"{uuid.uuid4().hex}_{resume.filename}"
            with path.open("wb") as fp:
                fp.write(content)
            quals = _extract_text_from_file(path) or quals
    if not quals:
        raise HTTPException(status_code=400, detail="Provide qualifications or upload a resume.")

    job_id = f"interview_{uuid.uuid4().hex}"
    questions = get_interview_questions(
        company_name=company_name,
        job_position=job_position,
        company_mission=company_mission or "Not provided",
        qualifications=quals,
    )
    if not questions:
        raise HTTPException(status_code=500, detail="Failed to generate questions.")

    job_record = {
        "status": "in_progress",
        "type": "interview",
        "created_at": datetime.utcnow().isoformat() + "Z",
        "job_info": {
            "company_name": company_name,
            "job_position": job_position,
            "company_mission": company_mission or "",
        },
        "qualifications": quals[:3000],
        "questions": questions,
        "answers": [],
    }
    _write_job(job_id, job_record)
    (INTERVIEW_DIR / job_id).mkdir(parents=True, exist_ok=True)
    return {"job_id": job_id, "questions": questions}


@app.post("/api/interview/job/{job_id}/answer")
async def submit_interview_answer(
    job_id: str,
    question_index: int = Form(...),
    audio_file: UploadFile = File(...),
):
    """Submit audio answer for a question. Transcribes and stores."""
    job = _read_job(job_id)
    if job.get("status") not in ("in_progress", "ready"):
        raise HTTPException(status_code=400, detail="Interview not in progress")
    questions = job.get("questions", [])
    if question_index < 0 or question_index >= len(questions):
        raise HTTPException(status_code=400, detail="Invalid question index")

    content = await audio_file.read()
    job_dir = INTERVIEW_DIR / job_id
    job_dir.mkdir(parents=True, exist_ok=True)
    ext = Path(audio_file.filename or "audio").suffix.lower() or ".webm"
    audio_path = job_dir / f"answer_{question_index}{ext}"
    with audio_path.open("wb") as fp:
        fp.write(content)

    wav_path = job_dir / f"answer_{question_index}.wav"
    try:
        audio.convert_to_wav(audio_path, wav_path)
        t = transcribe(wav_path)
        transcript_text = t.get("text", "").strip()
    except Exception:
        transcript_text = ""

    answers = job.get("answers", [])
    while len(answers) <= question_index:
        answers.append({"question": "", "transcript": "", "audio_id": ""})
    answers[question_index] = {
        "question": questions[question_index],
        "transcript": transcript_text,
        "audio_id": f"{job_id}_answer_{question_index}",
    }
    job["answers"] = answers
    _write_job(job_id, job)
    return {"question_index": question_index, "transcript": transcript_text}


@app.post("/api/interview/job/{job_id}/complete")
async def complete_interview(job_id: str, background_tasks: BackgroundTasks):
    """Mark interview complete and run Gemini analysis."""
    job = _read_job(job_id)
    if job.get("status") not in ("in_progress", "ready"):
        raise HTTPException(status_code=400, detail="Interview not in progress")
    job["status"] = "processing"
    _write_job(job_id, job)
    background_tasks.add_task(_process_interview_job, job_id)
    return {"job_id": job_id, "status": "processing"}


@app.get("/api/interview/job/{job_id}")
def get_interview_job(job_id: str):
    return _read_job(job_id)


@app.get("/api/interview/job/{job_id}/audio/{index}")
def get_interview_audio(job_id: str, index: int):
    """Serve recorded answer audio for playback."""
    job_dir = INTERVIEW_DIR / job_id
    wav_path = job_dir / f"answer_{index}.wav"
    if not wav_path.exists():
        raise HTTPException(status_code=404, detail="Audio not found")
    return FileResponse(wav_path, media_type="audio/wav")


def _process_interview_job(job_id: str) -> None:
    job = _read_job(job_id)
    try:
        ji = job.get("job_info", {})
        questions = job.get("questions", [])
        answers = job.get("answers", [])
        transcripts = [a.get("transcript", "") for a in answers if isinstance(a, dict)]
        while len(transcripts) < len(questions):
            transcripts.append("")

        analysis = get_interview_analysis(
            questions=questions,
            transcripts=transcripts,
            company_name=ji.get("company_name", ""),
            job_position=ji.get("job_position", ""),
            qualifications=job.get("qualifications", ""),
        )
        job.update(
            {
                "status": "done",
                "completed_at": datetime.utcnow().isoformat() + "Z",
                "result": {
                    "summary": analysis.get("summary", ""),
                    "score": analysis.get("score", 50),
                    "improvements": analysis.get("improvements", []),
                },
            }
        )
    except Exception as exc:
        job["status"] = "error"
        job["error"] = str(exc)
    finally:
        _write_job(job_id, job)


@app.get("/api/audio/{clip_id}")
def get_clip(clip_id: str):
    clip_path = _locate_file(CLIPS_DIR, clip_id)
    return FileResponse(clip_path, media_type="audio/wav")


@app.get("/api/tts/{tts_id}")
def get_tts(tts_id: str):
    tts_path = _locate_file(TTS_DIR, tts_id)
    return FileResponse(tts_path, media_type="audio/wav")


def _locate_file(root: Path, identifier: str) -> Path:
    candidate = root / f"{identifier}.wav"
    if candidate.exists():
        return candidate
    raise HTTPException(status_code=404, detail="Clip not found")


def _select_sample(sample: Optional[str]) -> Path:
    files = sorted(SAMPLES_DIR.glob("*.wav"))
    if not files:
        raise HTTPException(status_code=400, detail="No samples are available")
    if sample:
        requested = SAMPLES_DIR / f"{sample}.wav"
        if requested.exists():
            return requested
    return files[0]


def _process_job(job_id: str, upload_path: Path) -> None:
    job = _read_job(job_id)
    try:
        wav_path = audio.convert_to_wav(upload_path, UPLOADS_DIR / f"{job_id}.wav")
        playback_path = audio.convert_to_playback_wav(
            upload_path, UPLOADS_DIR / f"{job_id}_playback.wav"
        )
        transcript = transcribe(wav_path)
        aligned_words = align_words(
            wav_path,
            transcript.get("segments", []),
            language=transcript.get("language", "en"),
        )
        merged_words = _merge_alignment(transcript["words"], aligned_words)
        processed_words = _build_word_payload(job_id, wav_path, merged_words, playback_path=playback_path)
        summary = _summarize_words(processed_words)
        job.update(
            {
                "status": "done",
                "completed_at": datetime.utcnow().isoformat() + "Z",
                "result": {
                    "transcript": processed_words,
                    "summary": summary,
                    "text": transcript.get("text", ""),
                },
            }
        )
    except Exception as exc:  # pragma: no cover - failure path
        job["status"] = "error"
        job["error"] = str(exc)
    finally:
        _write_job(job_id, job)


def _build_word_payload(
    job_id: str,
    wav_path: Path,
    words: List[Dict],
    playback_path: Path | None = None,
) -> List[Dict]:
    processed: List[Dict] = []
    audio_segment = clip.load_audio(playback_path or wav_path)
    for index, entry in enumerate(words):
        scored = score.score_word(entry)
        prev_end = None
        next_start = None
        if index > 0:
            prev_end = words[index - 1].get("end")
        if index + 1 < len(words):
            next_start = words[index + 1].get("start")
        clip_id, clip_path = clip.slice_word_clip(
            playback_path or wav_path,
            scored["start"],
            scored["end"],
            CLIPS_DIR,
            job_id,
            prev_end=prev_end,
            next_start=next_start,
            audio_segment=audio_segment,
        )
        scored["clip_id"] = clip_id
        scored["observed_phonemes"] = recognize_phonemes(clip_path)
        if scored["issue"] != "ok":
            scored["tts_id"] = tts.synthesize_example(scored["word"], TTS_DIR)
        processed.append(scored)
    return processed


def _merge_alignment(words: List[Dict], aligned_words: List[Dict] | None) -> List[Dict]:
    if not aligned_words:
        return words
    aligned_index = 0
    merged: List[Dict] = []
    for word in words:
        normalized = _normalize_token(word.get("word", ""))
        match = None
        while aligned_index < len(aligned_words):
            candidate = aligned_words[aligned_index]
            aligned_index += 1
            if _normalize_token(candidate.get("word", "")) == normalized:
                match = candidate
                break
        if match:
            word["start"] = float(match.get("start", word.get("start", 0.0)))
            word["end"] = float(match.get("end", word.get("end", 0.0)))
            word["alignment_score"] = float(match.get("score", 0.0))
        merged.append(word)
    return merged


def _normalize_token(token: str) -> str:
    return "".join(ch for ch in token.lower() if ch.isalpha())


def _summarize_words(words: List[Dict]) -> Dict:
    if not words:
        return {"overall_score": 0, "top_issues": [], "practice": []}
    overall = round(sum(w["score"] for w in words) / len(words), 1)
    issue_counts: Dict[str, int] = {}
    for w in words:
        issue = w.get("issue")
        if issue and issue != "ok":
            issue_counts[issue] = issue_counts.get(issue, 0) + 1
    top_issues = sorted(issue_counts.items(), key=lambda item: item[1], reverse=True)[:3]
    practice = [w["word"] for w in words if w.get("issue") != "ok"]
    if len(practice) < 3:
        filler = sorted(words, key=lambda w: w["score"])[: 3 - len(practice)]
        practice.extend(w["word"] for w in filler)
    return {
        "overall_score": overall,
        "top_issues": [
            {"label": label, "count": count} for label, count in top_issues
        ],
        "practice": practice[:3],
    }


def _process_presentation_job(job_id: str, slides_path: Path, media_path: Path) -> None:
    job = _read_job(job_id)
    try:
        duration = _probe_duration(media_path)
        slide_count = _count_slides(slides_path)
        pacing = _estimate_pacing(duration, slide_count)
        base_tips = _presentation_tips(duration, slide_count)

        transcript_text = ""
        slide_texts: List[str] = []

        wav_path = _ensure_presentation_audio(media_path, job_id)
        if wav_path:
            try:
                transcript = transcribe(wav_path)
                transcript_text = transcript.get("text", "")
            except Exception:
                pass

        slide_texts = _extract_slide_text(slides_path)

        ai_tips = get_gemini_tips(
            transcript=transcript_text,
            slide_texts=slide_texts,
            pacing_label=pacing.get("label", "unknown"),
            slides_per_minute=pacing.get("slides_per_minute"),
            duration_seconds=duration,
        )
        audience_questions = get_audience_questions(transcript=transcript_text, slide_texts=slide_texts)

        result = {
            "slides": {
                "count": slide_count,
                "filename": slides_path.name,
            },
            "video": {
                "duration_seconds": duration,
                "filename": media_path.name,
                "media_type": "video" if media_path.suffix.lower() in {".mp4", ".mov", ".m4v", ".webm"} else "audio",
            },
            "coaching": {
                "pacing": pacing,
                "tips": (ai_tips or []) + base_tips,
                "ai_tips": ai_tips if ai_tips else None,
                "base_tips": base_tips,
            },
            "analysis": {
                "transcript": {"text": transcript_text} if transcript_text else None,
                "slides_text": slide_texts if slide_texts else None,
            },
            "audience_questions": audience_questions if audience_questions else None,
        }
        job.update(
            {
                "status": "done",
                "completed_at": datetime.utcnow().isoformat() + "Z",
                "result": result,
            }
        )
    except Exception as exc:  # pragma: no cover - failure path
        job["status"] = "error"
        job["error"] = str(exc)
    finally:
        _write_job(job_id, job)


def _ensure_presentation_audio(media_path: Path, job_id: str) -> Path | None:
    """Extract or convert media to 16k mono WAV for transcription. Returns path or None."""
    try:
        wav_path = PRESENTATION_DIR / job_id / f"{job_id}_audio.wav"
        wav_path.parent.mkdir(parents=True, exist_ok=True)
        audio.convert_to_wav(media_path, wav_path)
        return wav_path
    except Exception:
        return None


def _extract_slide_text(slides_path: Path) -> List[str]:
    """Extract text from slides (PPTX/PDF) or documents (DOCX). Returns empty list on failure."""
    suffix = slides_path.suffix.lower()
    if suffix in {".ppt", ".pptx"}:
        try:
            from pptx import Presentation

            pres = Presentation(str(slides_path))
            out = []
            for slide in pres.slides:
                lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = str(shape.text).strip()
                        if t:
                            lines.append(t)
                if lines:
                    out.append(" ".join(lines))
            return out
        except Exception:
            return []
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(slides_path))
            return [str(p.extract_text() or "").strip() for p in reader.pages]
        except Exception:
            return []
    if suffix == ".docx":
        return _extract_docx_text(slides_path)
    return []


def _extract_text_from_file(file_path: Path) -> str:
    """Extract full text from PDF, DOCX, or TXT. Returns empty string on failure."""
    suffix = file_path.suffix.lower()
    if suffix == ".txt":
        try:
            return file_path.read_text(encoding="utf-8", errors="replace").strip()
        except Exception:
            return ""
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            return " ".join(str(p.extract_text() or "").strip() for p in reader.pages).strip()
        except Exception:
            return ""
    if suffix == ".docx":
        try:
            from docx import Document

            doc = Document(str(file_path))
            return " ".join(p.text.strip() for p in doc.paragraphs if p.text.strip()).strip()
        except Exception:
            return ""
    return ""


def _extract_docx_text(docx_path: Path) -> List[str]:
    """Extract text from DOCX, chunked by sections/paragraphs for pacing."""
    try:
        from docx import Document

        doc = Document(str(docx_path))
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        if not lines:
            return []
        # Chunk into sections (~5 paragraphs each) for pacing
        chunk_size = max(1, len(lines) // 5)
        return [" ".join(lines[i : i + chunk_size]) for i in range(0, len(lines), chunk_size)]
    except Exception:
        return []


def _probe_duration(video_path: Path) -> float | None:
    import subprocess

    cmd = [
        "ffprobe",
        "-v",
        "error",
        "-show_entries",
        "format=duration",
        "-of",
        "default=noprint_wrappers=1:nokey=1",
        str(video_path),
    ]
    process = subprocess.run(cmd, capture_output=True, text=True)
    if process.returncode != 0:
        return None
    output = process.stdout.strip()
    try:
        return round(float(output), 2)
    except ValueError:
        return None


def _count_slides(slides_path: Path) -> int | None:
    suffix = slides_path.suffix.lower()
    if suffix in {".ppt", ".pptx"}:
        try:
            from pptx import Presentation

            presentation = Presentation(str(slides_path))
            return len(presentation.slides)
        except Exception:
            return None
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(slides_path))
            return len(reader.pages)
        except Exception:
            return None
    if suffix == ".docx":
        try:
            from docx import Document

            doc = Document(str(slides_path))
            para_count = sum(1 for p in doc.paragraphs if p.text.strip())
            return max(1, para_count // 5) if para_count else 1
        except Exception:
            return None
    return None


def _estimate_pacing(duration: float | None, slide_count: int | None) -> Dict:
    if not duration or not slide_count:
        return {
            "slides_per_minute": None,
            "label": "unknown",
        }
    minutes = duration / 60.0
    if minutes <= 0:
        return {"slides_per_minute": None, "label": "unknown"}
    rate = round(slide_count / minutes, 2)
    if rate < 0.8:
        label = "slow"
    elif rate <= 2.0:
        label = "balanced"
    else:
        label = "fast"
    return {"slides_per_minute": rate, "label": label}


def _presentation_tips(duration: float | None, slide_count: int | None) -> List[str]:
    tips = [
        "Keep your shoulders relaxed and weight balanced for confident posture.",
        "Use open hand gestures at key points to reinforce important ideas.",
        "Pause briefly after each section to let the audience absorb key points.",
    ]
    if duration and slide_count:
        pacing = _estimate_pacing(duration, slide_count)
        label = pacing.get("label")
        if label == "slow":
            tips.insert(0, "Your pacing feels slow—consider combining slides or tightening transitions.")
        elif label == "fast":
            tips.insert(0, "Your pacing feels fast—try pausing between slides to emphasize key points.")
        else:
            tips.insert(0, "Your pacing looks balanced—keep that steady rhythm.")
    return tips
