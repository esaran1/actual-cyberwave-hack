from __future__ import annotations

import math
import subprocess
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import numpy as np

from speech.transcribe import transcribe

try:  # pragma: no cover - optional heavy deps
    import cv2
    import mediapipe as mp
except Exception:  # pragma: no cover
    cv2 = None  # type: ignore
    mp = None  # type: ignore

try:  # pragma: no cover - optional heavy deps
    import librosa
except Exception:  # pragma: no cover
    librosa = None  # type: ignore


FILLER_WORDS = {"um", "uh", "like", "you know", "so", "actually", "basically"}


def analyze_presentation_media(media_path: Path, slides_path: Path | None = None) -> Dict:
    audio_path = _ensure_audio(media_path)
    audio_metrics, transcript_data = _analyze_audio(audio_path)
    video_metrics = _analyze_video(media_path)
    slide_text = _extract_slide_text(slides_path) if slides_path else []
    tips = _build_tips(audio_metrics, video_metrics)
    return {
        "audio": audio_metrics,
        "video": video_metrics,
        "transcript": transcript_data,
        "slides_text": slide_text,
        "tips": tips,
    }


def _ensure_audio(media_path: Path) -> Path:
    if media_path.suffix.lower() in {".wav", ".mp3", ".m4a", ".aac", ".flac", ".ogg"}:
        return media_path
    audio_path = media_path.with_suffix(".wav")
    cmd = ["ffmpeg", "-y", "-i", str(media_path), "-ac", "1", "-ar", "16000", str(audio_path)]
    subprocess.run(cmd, capture_output=True, text=True)
    return audio_path


def _analyze_audio(audio_path: Path) -> Tuple[Dict, Dict]:
    transcript = transcribe(audio_path)
    words = transcript.get("words", [])
    total_words = len(words)
    duration = _duration_from_words(words)
    minutes = duration / 60.0 if duration else None
    wpm = round(total_words / minutes, 1) if minutes and minutes > 0 else None

    pauses = _pause_stats(words)
    filler_stats = _filler_stats(words)

    metrics = {
        "duration_seconds": round(duration, 2) if duration else None,
        "word_count": total_words,
        "wpm": wpm,
        "pause_ratio": pauses["pause_ratio"],
        "long_pauses": pauses["long_pauses"],
        "filler_count": filler_stats["count"],
        "filler_rate": filler_stats["rate"],
    }
    return metrics, {
        "text": transcript.get("text", ""),
        "language": transcript.get("language", "en"),
    }


def _duration_from_words(words: List[Dict]) -> float:
    if not words:
        return 0.0
    end_times = [float(word.get("end") or 0.0) for word in words]
    return max(end_times) if end_times else 0.0


def _pause_stats(words: List[Dict]) -> Dict:
    if len(words) < 2:
        return {"pause_ratio": None, "long_pauses": 0}
    pauses = []
    for idx in range(len(words) - 1):
        end_time = float(words[idx].get("end") or 0.0)
        next_start = float(words[idx + 1].get("start") or end_time)
        pause = max(0.0, next_start - end_time)
        pauses.append(pause)
    long_pauses = len([p for p in pauses if p >= 0.8])
    total_pause = sum(pauses)
    speech_duration = _duration_from_words(words)
    pause_ratio = round(total_pause / speech_duration, 3) if speech_duration else None
    return {"pause_ratio": pause_ratio, "long_pauses": long_pauses}


def _filler_stats(words: List[Dict]) -> Dict:
    tokens = [str(word.get("word", "")).lower() for word in words]
    count = 0
    for token in tokens:
        if token in FILLER_WORDS:
            count += 1
    rate = round(count / max(1, len(tokens)), 3)
    return {"count": count, "rate": rate}


def _analyze_video(media_path: Path) -> Dict:
    if cv2 is None or mp is None:
        return {
            "supported": False,
            "gesture_rate": None,
            "posture_score": None,
            "gaze_variance": None,
        }

    capture = cv2.VideoCapture(str(media_path))
    if not capture.isOpened():
        return {
            "supported": False,
            "gesture_rate": None,
            "posture_score": None,
            "gaze_variance": None,
        }

    fps = capture.get(cv2.CAP_PROP_FPS) or 30
    sample_every = max(1, int(fps / 2))  # ~2 fps
    total_frames = int(capture.get(cv2.CAP_PROP_FRAME_COUNT)) or 0

    pose = mp.solutions.pose.Pose(static_image_mode=False)
    hands = mp.solutions.hands.Hands(static_image_mode=False, max_num_hands=2)

    last_wrists = None
    gesture_events = 0
    posture_scores = []
    gaze_offsets = []

    frame_idx = 0
    while True:
        ok, frame = capture.read()
        if not ok:
            break
        frame_idx += 1
        if frame_idx % sample_every != 0:
            continue

        rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
        pose_result = pose.process(rgb)
        hands_result = hands.process(rgb)

        if pose_result.pose_landmarks:
            landmarks = pose_result.pose_landmarks.landmark
            left_shoulder = landmarks[11]
            right_shoulder = landmarks[12]
            nose = landmarks[0]

            shoulder_width = math.dist(
                (left_shoulder.x, left_shoulder.y),
                (right_shoulder.x, right_shoulder.y),
            )
            shoulder_slope = abs(left_shoulder.y - right_shoulder.y)
            posture_score = max(0.0, 1.0 - (shoulder_slope / max(0.001, shoulder_width)))
            posture_scores.append(posture_score)

            shoulder_center_x = (left_shoulder.x + right_shoulder.x) / 2
            gaze_offsets.append(abs(nose.x - shoulder_center_x))

        wrists = []
        if hands_result.multi_hand_landmarks:
            for hand_landmarks in hands_result.multi_hand_landmarks:
                wrist = hand_landmarks.landmark[0]
                wrists.append((wrist.x, wrist.y))

        if wrists:
            if last_wrists:
                avg_move = 0.0
                for idx, wrist in enumerate(wrists):
                    if idx < len(last_wrists):
                        avg_move += math.dist(wrist, last_wrists[idx])
                avg_move /= max(1, len(wrists))
                if avg_move > 0.015:
                    gesture_events += 1
            last_wrists = wrists

    capture.release()
    pose.close()
    hands.close()

    duration_seconds = total_frames / fps if fps else None
    gesture_rate = round(gesture_events / max(1, duration_seconds / 60), 2) if duration_seconds else None
    posture_score = round(float(np.mean(posture_scores)), 3) if posture_scores else None
    gaze_variance = round(float(np.mean(gaze_offsets)), 3) if gaze_offsets else None

    return {
        "supported": True,
        "gesture_rate": gesture_rate,
        "posture_score": posture_score,
        "gaze_variance": gaze_variance,
    }


def _build_tips(audio_metrics: Dict, video_metrics: Dict) -> List[str]:
    tips = []
    wpm = audio_metrics.get("wpm")
    if wpm:
        if wpm > 170:
            tips.append("Your pace is fast. Try pausing after each key point.")
        elif wpm < 110:
            tips.append("Your pace is slow. Tighten transitions to keep energy up.")
        else:
            tips.append("Your pace is balanced. Keep the steady rhythm.")
    if audio_metrics.get("filler_rate") and audio_metrics["filler_rate"] > 0.05:
        tips.append("Reduce filler words by pausing silently instead of filling gaps.")
    if audio_metrics.get("pause_ratio") and audio_metrics["pause_ratio"] > 0.25:
        tips.append("Long pauses detected—use shorter pauses to maintain momentum.")

    if video_metrics.get("supported"):
        if video_metrics.get("gesture_rate") and video_metrics["gesture_rate"] < 6:
            tips.append("Add more hand gestures to emphasize key ideas.")
        if video_metrics.get("posture_score") and video_metrics["posture_score"] < 0.85:
            tips.append("Keep shoulders level to project confident posture.")
        if video_metrics.get("gaze_variance") and video_metrics["gaze_variance"] > 0.08:
            tips.append("Try keeping your gaze centered to stay connected with the audience.")
    else:
        tips.append("Video analysis unavailable—use a recorded video for posture and gesture feedback.")

    if not tips:
        tips.append("Great balance overall—keep practicing for even more polish.")
    return tips


def _extract_slide_text(slides_path: Path | None) -> List[str]:
    if slides_path is None:
        return []
    suffix = slides_path.suffix.lower()
    if suffix in {".ppt", ".pptx"}:
        try:
            from pptx import Presentation

            presentation = Presentation(str(slides_path))
            slide_text = []
            for slide in presentation.slides:
                lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = str(shape.text).strip()
                        if text:
                            lines.append(text)
                if lines:
                    slide_text.append(" ".join(lines))
            return slide_text
        except Exception:
            return []
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(slides_path))
            return [str(page.extract_text() or "").strip() for page in reader.pages]
        except Exception:
            return []
    return []
